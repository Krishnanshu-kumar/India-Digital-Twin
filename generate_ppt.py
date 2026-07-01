from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

def generate_ppt():
    template_path = "[Pub] ISRO BAH 2026 _ Idea Submission Template.pptx"
    output_path = "isro_pitch_deck_final.pptx"
    
    prs = Presentation(template_path)
    
    # --- Slide 1 ---
    slide = prs.slides[0]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if "Team Name :" in shape.text:
            shape.text = "Team Name : GeoInnovators"
        if "Problem Statement :" in shape.text:
            shape.text = "Problem Statement : Addressing the critical need for localized, high-resolution climate impact simulations to drive data-informed policy making and resilience planning across India."
        if "Team Leader Name :" in shape.text:
            shape.text = "Team Leader Name : Krishnanshu Kumar"

    # --- Slide 2: Team Members ---
    slide = prs.slides[1]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if "Team Members" in shape.text:
            p = shape.text_frame.add_paragraph()
            p.text = "1. Krishnanshu Kumar (Lead Architect)"
            p = shape.text_frame.add_paragraph()
            p.text = "2. [Teammate 2]"
            p = shape.text_frame.add_paragraph()
            p.text = "3. [Teammate 3]"
            p = shape.text_frame.add_paragraph()
            p.text = "4. [Teammate 4]"

    # --- Slide 3: Opportunity ---
    slide = prs.slides[2]
    for shape in slide.shapes:
        if shape.has_text_frame and "How different is it" in shape.text:
            shape.text = "How different is it from any of the other existing ideas?\nWe are moving beyond static historical data. Our Digital Twin provides interactive, advection-diffusion-based spatial simulations of the future (up to 2100) right in the browser.\n\nHow will it be able to solve the problem?\nIt empowers policymakers with real-time visual projections by combining IMD ground data, MOSDAC satellite data, and NICES ECV parameters into a unified dashboard.\n\nUSP of the proposed solution\nUltra-high resolution (0.25°) synthesis of ground and satellite data into a seamless, interactive simulation engine."

    # --- Slide 4: Features ---
    slide = prs.slides[3]
    for shape in slide.shapes:
        if shape.has_text_frame and "List of features" in shape.text:
            shape.text = "Key Features:\n• Interactive PyDeck Heatmaps and spatial anomaly visualizations.\n• Predictive Simulation Engine with seasonal adjustments and RCP Extreme scenarios.\n• Multi-Source Data Harmonization handling HDF5, NetCDF, and CSV formats natively.\n• Dynamic Time-Series charting and baseline comparisons for historical vs future data."

    # --- Slide 5: Process Flow ---
    slide = prs.slides[4]
    for shape in slide.shapes:
        if shape.has_text_frame and "Process flow diagram" in shape.text:
            shape.text = "Process Flow:\n1. Data Ingestion: Raw IMD, MOSDAC, and NICES datasets are downloaded.\n2. Harmonization: Data is resampled to a unified 0.25° grid and standardized.\n3. Simulation Engine: Future projections are generated using linear/exponential extrapolation and advection physics.\n4. UI Layer: Streamlit renders the data interactively using Pandas and PyDeck."

    # --- Slide 6: Wireframes ---
    slide = prs.slides[5]
    for shape in slide.shapes:
        if shape.has_text_frame and "Wireframes" in shape.text:
            shape.text = "Wireframes / Screenshots (See Attached Image)"
            # Insert the map screenshot we have saved in brain folder
            img_path = r"C:\Users\KIIT\.gemini\antigravity-ide\brain\42880b52-7dc6-4a32-9a1a-2fa479697fa7\surface_ocean_maps_1782818578385.png"
            if Path(img_path).exists():
                slide.shapes.add_picture(img_path, Inches(1), Inches(2), height=Inches(4.5))

    # --- Slide 7: Architecture ---
    slide = prs.slides[6]
    for shape in slide.shapes:
        if shape.has_text_frame and "Architecture diagram" in shape.text:
            shape.text = "Architecture Summary:\n• Frontend: Streamlit, Deck.gl, Plotly\n• Computation Backend: NumPy, SciPy (Gaussian smoothing), Xarray\n• Data Layer: Pandas DataFrames for quick spatial rendering\n• Orchestration: Native Python multithreading for background downloads"

    # --- Slide 8: Technologies ---
    slide = prs.slides[7]
    for shape in slide.shapes:
        if shape.has_text_frame and "Technologies to be used" in shape.text:
            shape.text = "Technologies Used:\n• Programming Language: Python 3\n• UI Framework: Streamlit\n• Spatial Visualization: PyDeck (deck.gl)\n• Data Processing: Pandas, Xarray, H5py\n• Math & Physics: NumPy, SciPy"

    # --- Slide 9: Cost ---
    slide = prs.slides[8]
    for shape in slide.shapes:
        if shape.has_text_frame and "Estimated implementation cost" in shape.text:
            shape.text = "Estimated Implementation Cost:\n• Development & Tooling: $0 (100% Open Source Python Stack)\n• Cloud Hosting (App Service): ~$20 - $50 / month\n• Data Storage (AWS S3 or equivalent for large NetCDF files): ~$100 / month depending on historical data volume.\n• Total MVP Cost: <$150 / month"

    prs.save(output_path)
    print(f"Successfully generated {output_path}")

if __name__ == "__main__":
    generate_ppt()
